from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO, MSO_SHAPE, MSO_CONNECTOR

prs = Presentation('./ppt_file/eg2.pptx')

text_runs = []
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

print(text_runs)

## add_picture()
